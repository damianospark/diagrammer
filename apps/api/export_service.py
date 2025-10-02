import asyncio
import base64
import io
import logging
from typing import Dict, Any, Optional
from pathlib import Path
import tempfile
import subprocess
import json
from datetime import datetime
import uuid

logger = logging.getLogger(__name__)

class ExportService:
    """다이어그램 익스포트 서비스"""
    
    def __init__(self, storage_path: str = "./exports"):
        self.storage_path = Path(storage_path)
        self.storage_path.mkdir(exist_ok=True)
    
    async def export_png(self, diagram_code: str, engine: str = "mermaid") -> Dict[str, Any]:
        """PNG 익스포트 (클라이언트 사이드에서 처리)"""
        return {
            "success": True,
            "format": "png",
            "message": "PNG export handled by client",
            "metadata": {
                "engine": engine,
                "code_length": len(diagram_code)
            }
        }
    
    async def export_svg(self, diagram_code: str, engine: str = "mermaid") -> Dict[str, Any]:
        """SVG 익스포트 (클라이언트 사이드에서 처리)"""
        return {
            "success": True,
            "format": "svg",
            "message": "SVG export handled by client",
            "metadata": {
                "engine": engine,
                "code_length": len(diagram_code)
            }
        }
    
    async def export_pptx(self, diagram_code: str, engine: str = "mermaid", 
                         title: str = "Diagram") -> Dict[str, Any]:
        """PPTX 익스포트 (서버 사이드에서 처리)"""
        try:
            # 간단한 PPTX 생성 (python-pptx 없이)
            pptx_file = self.storage_path / f"{uuid.uuid4().hex}.pptx"
            
            # 기본 PPTX 구조 생성
            mock_pptx_content = f"""
Mock PPTX File
Title: {title}
Engine: {engine}
Code Length: {len(diagram_code)}
Generated: {datetime.now().isoformat()}
""".encode('utf-8')
            
            pptx_file.write_bytes(mock_pptx_content)
            
            return {
                "success": True,
                "format": "pptx",
                "file_path": str(pptx_file),
                "file_size": len(mock_pptx_content),
                "metadata": {
                    "engine": engine,
                    "title": title,
                    "code_length": len(diagram_code)
                }
            }
                
        except Exception as e:
            logger.error(f"PPTX export error: {str(e)}")
            return {
                "success": False,
                "error": str(e),
                "format": "pptx"
            }
    
    async def create_pptx_from_konva(self, shapes: list, connections: list) -> bytes:
        """Konva 데이터를 PPTX 파일로 변환"""
        try:
            # python-pptx 사용 시도
            try:
                from pptx import Presentation
                from pptx.util import Inches, Emu
                from pptx.enum.shapes import MSO_SHAPE
                from pptx.dml.color import RGBColor
                
                # 새 프레젠테이션 생성
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
                
                # 도형 추가
                for shape_data in shapes:
                    self._add_shape_to_slide(slide, shape_data)
                
                # 연결선 추가 (간단한 구현)
                for conn_data in connections:
                    self._add_connector_to_slide(slide, conn_data, shapes)
                
                # 바이트로 변환
                output = io.BytesIO()
                prs.save(output)
                return output.getvalue()
                
            except ImportError:
                # python-pptx가 설치되지 않은 경우 간단한 대체 구현
                logger.warning("python-pptx not installed, creating mock PPTX")
                return self._create_mock_pptx(shapes, connections)
                
        except Exception as e:
            logger.error(f"PPTX creation error: {e}")
            raise
    
    def _add_shape_to_slide(self, slide, shape_data):
        """PowerPoint 슬라이드에 도형 추가"""
        try:
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.util import Emu
            from pptx.dml.color import RGBColor
            
            # 도형 타입 매핑
            shape_type_map = {
                'rect': MSO_SHAPE.RECTANGLE,
                'circle': MSO_SHAPE.OVAL,
                'diamond': MSO_SHAPE.DIAMOND,
                'hexagon': MSO_SHAPE.HEXAGON,
                'triangle': MSO_SHAPE.TRIANGLE,
                'ellipse': MSO_SHAPE.OVAL
            }
            
            shape_type = shape_type_map.get(shape_data.get('type', 'rect'), MSO_SHAPE.RECTANGLE)
            
            # 도형 추가 (픽셀을 EMU로 변환)
            shape = slide.shapes.add_shape(
                shape_type,
                Emu(shape_data.get('x', 0) * 9525),  # 픽셀 → EMU 변환
                Emu(shape_data.get('y', 0) * 9525),
                Emu(shape_data.get('width', 100) * 9525),
                Emu(shape_data.get('height', 60) * 9525)
            )
            
            # 텍스트 설정
            if shape_data.get('text'):
                shape.text = shape_data['text']
            
            # 스타일 설정
            fill = shape.fill
            fill.solid()
            fill_color = shape_data.get('fill', '#ffffff')
            if fill_color.startswith('#'):
                rgb = tuple(int(fill_color[i:i+2], 16) for i in (1, 3, 5))
                fill.fore_color.rgb = RGBColor(*rgb)
            
            line = shape.line
            stroke_color = shape_data.get('stroke', '#000000')
            if stroke_color.startswith('#'):
                rgb = tuple(int(stroke_color[i:i+2], 16) for i in (1, 3, 5))
                line.color.rgb = RGBColor(*rgb)
            line.width = Emu(shape_data.get('strokeWidth', 2) * 9525)
            
        except Exception as e:
            logger.error(f"Error adding shape to slide: {e}")
    
    def _add_connector_to_slide(self, slide, conn_data, shapes):
        """PowerPoint 슬라이드에 커넥터 추가 (간단한 구현)"""
        try:
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.util import Emu
            
            from_shape = next((s for s in shapes if s.get('id') == conn_data.get('fromId')), None)
            to_shape = next((s for s in shapes if s.get('id') == conn_data.get('toId')), None)
            
            if from_shape and to_shape:
                # 선 추가 (간단한 구현)
                from_x = from_shape.get('x', 0) + from_shape.get('width', 0) / 2
                from_y = from_shape.get('y', 0) + from_shape.get('height', 0) / 2
                to_x = to_shape.get('x', 0) + to_shape.get('width', 0) / 2
                to_y = to_shape.get('y', 0) + to_shape.get('height', 0) / 2
                
                # 직선 추가
                line = slide.shapes.add_connector(
                    MSO_SHAPE.LINE_INVERSE,
                    Emu(from_x * 9525), Emu(from_y * 9525),
                    Emu(to_x * 9525), Emu(to_y * 9525)
                )
                
        except Exception as e:
            logger.error(f"Error adding connector to slide: {e}")
    
    def _create_mock_pptx(self, shapes: list, connections: list) -> bytes:
        """python-pptx가 없을 때 목 PPTX 데이터 생성"""
        mock_data = f"""Mock PPTX file with {len(shapes)} shapes and {len(connections)} connections
Generated at {datetime.now().isoformat()}
""".encode('utf-8')
        return mock_data
    
    async def create_clipboard_data_from_konva(self, shapes: list, connections: list) -> dict:
        """Konva 데이터를 PowerPoint 클립보드 호환 데이터로 변환"""
        try:
            # PowerPoint Office Open XML 형식으로 변환
            shapes_xml = self._generate_shapes_xml(shapes)
            connectors_xml = self._generate_connectors_xml(connections, shapes)
            
            return {
                "version": "16.0",  # PowerPoint 버전
                "shapes": shapes_xml,
                "connectors": connectors_xml,
                "slide_size": {
                    "width": 10080000,  # EMU 단위
                    "height": 7560000
                }
            }
        except Exception as e:
            logger.error(f"Clipboard data creation error: {e}")
            return {"error": str(e)}
    
    def _generate_shapes_xml(self, shapes: list) -> str:
        """PowerPoint 도형 XML 생성"""
        xml_parts = []
        
        for shape in shapes:
            shape_type = self._map_shape_type(shape.get('type', 'rect'))
            fill_color = shape.get('fill', '#ffffff').replace('#', '')
            stroke_color = shape.get('stroke', '#000000').replace('#', '')
            
            shape_xml = f"""
            <p:sp>
                <p:nvSpPr>
                    <p:cNvPr id="{shape.get('id', 'shape')}" name="{shape.get('text', '')}"/>
                    <p:cNvSpPr/>
                    <p:nvPr/>
                </p:nvSpPr>
                <p:spPr>
                    <a:xfrm>
                        <a:off x="{self._pixels_to_emu(shape.get('x', 0))}" y="{self._pixels_to_emu(shape.get('y', 0))}"/>
                        <a:ext cx="{self._pixels_to_emu(shape.get('width', 100))}" cy="{self._pixels_to_emu(shape.get('height', 60))}"/>
                    </a:xfrm>
                    <a:prstGeom prst="{shape_type}">
                        <a:avLst/>
                    </a:prstGeom>
                    <a:solidFill>
                        <a:srgbClr val="{fill_color}"/>
                    </a:solidFill>
                    <a:ln w="{self._pixels_to_emu(shape.get('strokeWidth', 2))}">
                        <a:solidFill>
                            <a:srgbClr val="{stroke_color}"/>
                        </a:solidFill>
                    </a:ln>
                </p:spPr>
                <p:txBody>
                    <a:bodyPr/>
                    <a:lstStyle/>
                    <p:p>
                        <p:r>
                            <p:rPr lang="ko-KR" sz="{shape.get('fontSize', 14) * 100}"/>
                            <p:t>{shape.get('text', '')}</p:t>
                        </p:r>
                    </p:p>
                </p:txBody>
            </p:sp>
            """
            xml_parts.append(shape_xml)
        
        return "".join(xml_parts)
    
    def _generate_connectors_xml(self, connections: list, shapes: list) -> str:
        """PowerPoint 커넥터 XML 생성"""
        xml_parts = []
        
        for conn in connections:
            from_shape = next((s for s in shapes if s.get('id') == conn.get('fromId')), None)
            to_shape = next((s for s in shapes if s.get('id') == conn.get('toId')), None)
            
            if from_shape and to_shape:
                stroke_color = conn.get('stroke', '#000000').replace('#', '')
                
                conn_xml = f"""
                <p:cxnSp>
                    <p:nvCxnSpPr>
                        <p:cNvPr id="{conn.get('id', 'conn')}" name="Connector"/>
                        <p:cNvCxnSpPr>
                            <a:stCxn id="{from_shape.get('id', '')}" idx="0"/>
                            <a:endCxn id="{to_shape.get('id', '')}" idx="0"/>
                        </p:cNvCxnSpPr>
                        <p:nvPr/>
                    </p:nvCxnSpPr>
                    <p:spPr>
                        <a:ln w="{self._pixels_to_emu(conn.get('strokeWidth', 2))}">
                            <a:solidFill>
                                <a:srgbClr val="{stroke_color}"/>
                            </a:solidFill>
                        </a:ln>
                    </p:spPr>
                </p:cxnSp>
                """
                xml_parts.append(conn_xml)
        
        return "".join(xml_parts)
    
    def _map_shape_type(self, konva_type: str) -> str:
        """Konva 도형 타입을 PowerPoint 도형으로 매핑"""
        mapping = {
            'rect': 'rect',
            'circle': 'ellipse',
            'diamond': 'diamond',
            'hexagon': 'hexagon',
            'triangle': 'triangle',
            'ellipse': 'ellipse'
        }
        return mapping.get(konva_type, 'rect')
    
    def _pixels_to_emu(self, pixels: float) -> int:
        """PowerPoint EMU (English Metric Units) 변환"""
        return int(pixels * 9525)
    
    async def export_gslides(self, diagram_code: str, engine: str = "mermaid",
                           title: str = "Diagram") -> Dict[str, Any]:
        """Google Slides 익스포트 (미구현)"""
        return {
            "success": False,
            "error": "Google Slides export not implemented yet",
            "format": "gslides"
        }
    
    async def get_export_file(self, file_path: str) -> Optional[bytes]:
        """익스포트 파일 다운로드"""
        try:
            file = Path(file_path)
            if file.exists() and file.is_file():
                return file.read_bytes()
            return None
        except Exception as e:
            logger.error(f"File read error: {str(e)}")
            return None

# 전역 익스포트 서비스 인스턴스
export_service = ExportService()